"""Executive presentation generation: PPTX (python-pptx) + reveal.js HTML.

Slides are built deterministically from the *verified* report — every number on
a slide already passed the verification gate, so no LLM call (and no new
hallucination surface) is introduced here.
"""
from __future__ import annotations

import html as html_mod
import io
import logging
from pathlib import Path

from api.models.report import Report

logger = logging.getLogger(__name__)

INDIGO = (0x4F, 0x46, 0xE5)
DARKINK = (0x11, 0x18, 0x27)
GRAY = (0x6B, 0x72, 0x80)


def _slides_content(report: Report, workspace: Path) -> list[dict]:
    """Distill the report into 6-8 slide definitions (title, bullets, chart)."""
    chart_png = {c.name: workspace / c.png_path for c in report.charts if c.png_path}

    def png_for(name: str | None):
        p = chart_png.get(name or "")
        return p if p is not None and p.exists() else None

    slides: list[dict] = [{
        "kind": "title",
        "title": report.title,
        "subtitle": report.executive_summary[:300],
    }]
    if report.kpis:
        slides.append({
            "kind": "kpis", "title": {"en": "Key Metrics", "ar": "المؤشرات الرئيسية"}[report.language],
            "kpis": [(k.label, k.value, k.change) for k in report.kpis[:8]],
        })
    for f in report.findings[:3]:
        slides.append({
            "kind": "finding", "title": f.title,
            "bullets": [s.strip() for s in f.narrative.split(". ") if s.strip()][:4],
            "image": png_for(f.chart_name),
        })
    if report.segments:
        slides.append({
            "kind": "bullets", "title": {"en": "Segments", "ar": "الشرائح"}[report.language],
            "bullets": [f"{s.name}: {s.recommendation}" for s in report.segments[:5]],
        })
    if report.anomalies:
        slides.append({
            "kind": "bullets", "title": {"en": "Anomalies", "ar": "الحالات الشاذة"}[report.language],
            "bullets": [f"{a.title} — {a.narrative}"[:220] for a in report.anomalies[:3]],
        })
    if report.forecast:
        slides.append({
            "kind": "finding", "title": {"en": "Forecast", "ar": "التنبؤ"}[report.language],
            "bullets": [report.forecast.narrative[:220], report.forecast.reliability_statement[:220]],
            "image": png_for(report.forecast.chart_name),
        })
    if report.recommendations:
        slides.append({
            "kind": "bullets", "title": {"en": "Recommendations", "ar": "التوصيات"}[report.language],
            "bullets": report.recommendations[:6],
        })
    slides.append({
        "kind": "bullets", "title": {"en": "Appendix — Data Quality", "ar": "الملحق — جودة البيانات"}[report.language],
        "bullets": ([report.data_quality_notes] if report.data_quality_notes else [])
        + [f"{e.action}: {e.justification}"[:180] for e in report.cleaning_log[:4]],
    })
    return slides


# ----------------------------------------------------------------- PPTX ----

def render_pptx(report: Report, workspace: Path) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_text(slide, left, top, width, height, text, size, bold=False, color=DARKINK):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return box

    for s in _slides_content(report, workspace):
        slide = prs.slides.add_slide(blank)
        if s["kind"] == "title":
            add_text(slide, 0.8, 2.3, 11.7, 1.4, s["title"], 40, bold=True, color=INDIGO)
            add_text(slide, 0.8, 3.8, 11.7, 2.4, s.get("subtitle", ""), 16, color=GRAY)
            continue
        add_text(slide, 0.8, 0.5, 11.7, 0.9, str(s["title"]), 28, bold=True, color=INDIGO)
        if s["kind"] == "kpis":
            for i, (label, value, change) in enumerate(s["kpis"]):
                col, row = i % 4, i // 4
                x, y = 0.8 + col * 3.05, 1.8 + row * 2.3
                add_text(slide, x, y, 2.8, 0.8, str(value), 26, bold=True, color=INDIGO)
                add_text(slide, x, y + 0.7, 2.8, 0.5, str(label), 12, color=GRAY)
                if change:
                    add_text(slide, x, y + 1.1, 2.8, 0.5, str(change), 12, color=DARKINK)
        else:
            bullets = s.get("bullets") or []
            text_width = 6.0 if s.get("image") else 11.7
            for j, b in enumerate(bullets[:6]):
                add_text(slide, 0.8, 1.7 + j * 0.95, text_width, 0.9, f"•  {b}", 14)
            img = s.get("image")
            if img:
                try:
                    slide.shapes.add_picture(str(img), Inches(7.0), Inches(1.6),
                                             width=Inches(5.6))
                except Exception:  # noqa: BLE001 — slide remains text-only
                    logger.exception("could not embed %s", img)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ------------------------------------------------------------- reveal.js ----

def render_slides_html(report: Report, workspace: Path) -> str:
    rtl = report.language == "ar"
    sections = []
    for s in _slides_content(report, workspace):
        esc = html_mod.escape
        if s["kind"] == "title":
            sections.append(f"<section><h1>{esc(s['title'])}</h1><p class='sub'>{esc(s.get('subtitle',''))}</p></section>")
        elif s["kind"] == "kpis":
            cells = "".join(
                f"<div class='kpi'><div class='v'>{esc(str(v))}</div><div class='l'>{esc(str(l))}</div>"
                f"{f'<div class=c>{esc(str(ch))}</div>' if ch else ''}</div>"
                for l, v, ch in s["kpis"])
            sections.append(f"<section><h2>{esc(str(s['title']))}</h2><div class='kpis'>{cells}</div></section>")
        else:
            bullets = "".join(f"<li>{esc(str(b))}</li>" for b in (s.get("bullets") or []))
            img = ""
            p = s.get("image")
            if p:
                import base64
                img = (f"<img src='data:image/png;base64,"
                       f"{base64.b64encode(Path(p).read_bytes()).decode()}'/>")
            sections.append(f"<section><h2>{esc(str(s['title']))}</h2>"
                            f"<div class='split'><ul>{bullets}</ul>{img}</div></section>")

    return f"""<!DOCTYPE html>
<html lang="{report.language}" dir="{'rtl' if rtl else 'ltr'}">
<head><meta charset="utf-8"><title>{html_mod.escape(report.title)}</title>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reveal.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/theme/white.css">
<style>
.reveal h1, .reveal h2 {{ color:#4F46E5; font-family: Inter, 'Noto Sans Arabic', sans-serif; }}
.reveal .sub {{ color:#6B7280; font-size:.6em; }}
.reveal .kpis {{ display:grid; grid-template-columns:repeat(4,1fr); gap:18px; }}
.reveal .kpi {{ background:#EEF2FF; border-radius:12px; padding:18px; }}
.reveal .kpi .v {{ font-size:.9em; font-weight:700; color:#4F46E5; }}
.reveal .kpi .l {{ font-size:.4em; color:#6B7280; }}
.reveal .kpi .c {{ font-size:.45em; color:#111827; }}
.reveal .split {{ display:flex; gap:24px; align-items:center; }}
.reveal .split ul {{ font-size:.6em; flex:1; }}
.reveal .split img {{ max-width:46%; border-radius:8px; }}
</style></head>
<body><div class="reveal"><div class="slides">{''.join(sections)}</div></div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reveal.js"></script>
<script>Reveal.initialize({{rtl: {str(rtl).lower()}}});</script>
</body></html>"""
